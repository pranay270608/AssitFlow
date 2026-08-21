from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_assistflow_ppt():
    prs = Presentation()

    # Define NYC Palette Colors
    asphalt_black = RGBColor(18, 18, 18)
    subway_white = RGBColor(248, 250, 252)
    taxi_yellow = RGBColor(255, 204, 0)

    # Slide Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    image_slide_layout = prs.slide_layouts[5] # Title only for UI screenshots

    slides_data = [
        {
            "type": "title",
            "title": "AssistFlow",
            "content": "Enterprise AI Support & Operations Portal\nNext-Generation Automated IT Troubleshooting & Live Telemetry\n\nPresented by: Pranay"
        },
        {
            "type": "bullet",
            "title": "2. Introduction",
            "content": [
                "What is Support Ticket Management? The operational backbone tracking and resolving enterprise technical issues.",
                "What Does AssistFlow Do? An AI-first, centralized desk that intercepts and resolves tickets instantly using local LLMs.",
                "Why Use This System? Eliminate L1 bottlenecks, guarantee data privacy, and deliver 24/7 remediation globally."
            ]
        },
        {
            "type": "bullet",
            "title": "3. Problem Statement",
            "content": [
                "The Reactive Bottleneck: Employees wait hours for basic triage.",
                "Wasted Engineering Hours: Staff spend too much time categorizing simple tickets.",
                "Human Error: Critical outages get lost in high-volume, low-priority queues.",
                "High MTTR: Prolonged downtime directly impacts enterprise revenue."
            ]
        },
        {
            "type": "bullet",
            "title": "4. Existing System",
            "content": [
                "Current Handling: Fragmented emails and static intranet forms.",
                "The 'Black Hole' Effect: Zero immediate feedback for users.",
                "Manual Dependency: Every ticket requires a human reader.",
                "Siloed Knowledge: Resolutions depend on individual technician experience."
            ]
        },
        {
            "type": "bullet",
            "title": "5. Proposed Solution",
            "content": [
                "Our AI Approach: Intercepts tickets instantly and deploys a local AI engine to diagnose.",
                "Instant Diagnostics: Deflects simple tickets away from human staff with immediate solutions.",
                "Smart Escalation: Automatically formats and routes complex issues to the correct department.",
                "Enterprise Privacy: Localized LLMs ensure corporate infrastructure data never leaves the MNC."
            ]
        },
        {
            "type": "bullet",
            "title": "6. Key Objectives & Features",
            "content": [
                "AI Classification: Instantly categorizes Hardware, Software, Network, and Access.",
                "Priority Detection: Dynamic urgency analysis to flag Critical issues.",
                "Automated Assignment: Matches tickets to the right available technician.",
                "AI Diagnostic Chatbot: Virtual assistant for step-by-step troubleshooting.",
                "Live Dashboard: Real-time telemetry tracking IT health.",
                "Integrated Bug Routing: Parallel pipeline for software defects."
            ]
        },
        {
            "type": "bullet",
            "title": "7. System Architecture",
            "content": [
                "User Interface -> Streamlit Frontend -> Python Backend -> Ollama AI Engine -> MongoDB -> Technician",
                "Core Flow: Secure, bidirectional communication giving AI real-time database context."
            ]
        },
        {
            "type": "bullet",
            "title": "8. Technology Stack",
            "content": [
                "Frontend: Streamlit, Custom HTML/CSS, Plotly.",
                "Backend: Python 3, FastAPI / Flask, Secure Auth.",
                "Database: MongoDB (NoSQL high-speed document architecture).",
                "AI/ML Engine: Localized LLM APIs (Ollama, Llama 3, Mistral).",
                "Deployment: Docker-ready for internal enterprise hosting."
            ]
        },
        {
            "type": "bullet",
            "title": "9. AI/ML Workflow",
            "content": [
                "Input: 'GlobalProtect VPN drops every 5 minutes.'",
                "NLP Processing: Extracts key technical entities.",
                "Category Prediction: Maps to 'Network'.",
                "Priority & Sentiment: Analyzes urgency and flags as 'High'.",
                "Assignment & Remediation: Suggests fix or routes to Network Admin."
            ]
        },
        {
            "type": "image",
            "title": "10. Lifecycle & UI Showcase",
            "content": "INSERT YOUR CUSTOM ASSISTFLOW SCREENSHOTS HERE:\n- Login & Auth\n- Ticket Creation\n- AI Recommendations\n- Dashboard"
        },
        {
            "type": "bullet",
            "title": "11. Impact & Why Choose AssistFlow?",
            "content": [
                "Deflection Rate: Solves standard IT issues without human intervention.",
                "Strategic Reallocation: Technicians focus on high-level infrastructure engineering.",
                "Zero Cloud Dependency: Absolute data sovereignty with no 3rd-party API costs."
            ]
        },
        {
            "type": "bullet",
            "title": "12. Results & Benefits",
            "content": [
                "Faster Response: Seconds for AI triage vs. hours for manual.",
                "Automated Routing: Near 100% immediate categorization.",
                "AI Classification Accuracy: Consistent, bias-free urgency flagging.",
                "Reduced Workload: Reclaims hundreds of engineering hours."
            ]
        },
        {
            "type": "bullet",
            "title": "13. Future Scope & Conclusion",
            "content": [
                "Voice Support: Hands-free IT support logging.",
                "Multilingual AI: On-the-fly translation for global branches.",
                "Predictive Analytics: Identifying hardware trends before breakdowns.",
                "Auto-Resolution: AI executing background scripts for routine fixes.",
                "Conclusion: Transitioning IT from a reactive cost-center to a proactive powerhouse."
            ]
        },
        {
            "type": "title",
            "title": "14. Any Queries?",
            "content": "Open for Q&A"
        },
        {
            "type": "title",
            "title": "15. Thank You",
            "content": "AssistFlow\nEmpowering the future of enterprise operations.\n\nPresented By: Pranay"
        }
    ]

    for slide_data in slides_data:
        if slide_data["type"] == "title":
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            subtitle.text = slide_data["content"]
            
        elif slide_data["type"] == "bullet":
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            body = slide.shapes.placeholders[1]
            title.text = slide_data["title"]
            tf = body.text_frame
            for point in slide_data["content"]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
        elif slide_data["type"] == "image":
            slide = prs.slides.add_slide(image_slide_layout)
            title = slide.shapes.title
            title.text = slide_data["title"]
            # Placeholder text box for instructions
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
            tf = txBox.text_frame
            tf.text = slide_data["content"]

    prs.save('AssistFlow_Pitch_Deck.pptx')
    print("Success! Your presentation has been saved as 'AssistFlow_Pitch_Deck.pptx'")

if __name__ == "__main__":
    create_assistflow_ppt()